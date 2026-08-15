#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
OfficeConverter 纯 Python 后端（T7）单元测试

验证 ``_convert_pure_python`` 及四个降级后端（docx→HTML、pptx→大纲、
xlsx→表格、doc/xls/ppt→安装提示）在 LO / COM 均不可用时的行为：

- 6 种后缀在全部后端不可用时各自的 content_type / backend_used / 消息（Metis E6）
- xlsx 行 / 列上限（5000×200）与「已截断」标记（Metis E7）
- 单一库缺失时仅对应后缀降级失败，其余后缀仍正常工作（Metis B4-4）
- 文件缺失 / 损坏时返回结果而不抛出
- 惰性导入：模块加载不引入 mammoth / openpyxl / pptx

所有真实转换均使用 tmp_path 下由对应库现场生成的迷你夹具，不依赖外部文件。
"""

import builtins
import importlib
import sys
from contextlib import ExitStack
from unittest.mock import patch

import pytest

from freeassetfilter.services.office_converter import (
    ERROR_MESSAGE,
    ConversionResult,
    OfficeConverter,
)

# 模块加载不得触碰的第三方库（Metis B4-4 惰性导入验证目标）。
_LAZY_LIBS = ("mammoth", "openpyxl", "pptx")


# ===========================================================================
# 夹具：现场生成真实的 docx / pptx / xlsx 迷你文件
# ===========================================================================


@pytest.fixture
def office_file(tmp_path):
    """按后缀生成真实迷你文件，返回路径字符串。"""

    def _make(suffix: str) -> str:
        if suffix == "docx":
            from docx import Document

            doc = Document()
            doc.add_paragraph("Hello 世界")
            doc.add_paragraph("Second paragraph")
            path = tmp_path / "sample.docx"
            doc.save(path)
            return str(path)

        if suffix == "pptx":
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()
            blank = prs.slide_layouts[6]
            for title_text in ("Slide One", "Slide Two"):
                slide = prs.slides.add_slide(blank)
                textbox = slide.shapes.add_textbox(
                    Inches(1), Inches(1), Inches(4), Inches(1)
                )
                textbox.text = title_text
            path = tmp_path / "sample.pptx"
            prs.save(path)
            return str(path)

        if suffix == "xlsx":
            import openpyxl

            wb = openpyxl.Workbook()
            ws = wb.active
            ws.append(["名称", "数量"])
            ws.append(["苹果", 3])
            ws.append(["香蕉", 5])
            path = tmp_path / "sample.xlsx"
            wb.save(path)
            return str(path)

        # legacy：不会真正读取，任意字节即可。
        path = tmp_path / f"sample.{suffix}"
        path.write_bytes(b"dummy")
        return str(path)

    return _make


@pytest.fixture
def large_xlsx_path(tmp_path):
    """5200 行（无表头）的工作簿，用于验证 5000 行上限。"""
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    for i in range(5200):
        ws.append([f"row{i}", i])
    path = tmp_path / "large.xlsx"
    wb.save(path)
    return str(path)


@pytest.fixture
def wide_xlsx_path(tmp_path):
    """250 列的工作簿，用于验证 200 列上限。"""
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append([f"col{i}" for i in range(250)])
    path = tmp_path / "wide.xlsx"
    wb.save(path)
    return str(path)


@pytest.fixture
def all_backends_down():
    """LO 与 COM 探测均不可用（将现代格式路由到纯 Python）。"""
    with ExitStack() as stack:
        stack.enter_context(
            patch.object(OfficeConverter, "_soffice_available", return_value=False)
        )
        stack.enter_context(
            patch.object(OfficeConverter, "_com_available", return_value=False)
        )
        yield stack.pop_all()


# ===========================================================================
# Metis E6：6 种后缀在全部后端不可用时的降级行为
# ===========================================================================


class TestAllSuffixesAllBackendsDown:
    """LO / COM 均不可用时，6 种后缀各自应路由到正确的后端并给出正确类型。"""

    @pytest.mark.parametrize(
        "suffix, expected_type, expected_backend",
        [
            ("docx", "html", "pure-python"),
            ("pptx", "outline", "pure-python"),
            ("xlsx", "table", "pure-python"),
            ("doc", "error", "error"),
            ("xls", "error", "error"),
            ("ppt", "error", "error"),
        ],
    )
    def test_all_backends_unavailable_six_suffixes(
        self,
        suffix: str,
        expected_type: str,
        expected_backend: str,
        office_file,
        all_backends_down,
    ) -> None:
        """6 种后缀在 LO/COM 均不可用时的 content_type / backend_used / 消息。"""
        path = office_file(suffix)
        with all_backends_down:
            result = OfficeConverter.convert(
                {"path": path, "suffix": suffix, "is_dir": False}
            )

        assert isinstance(result, ConversionResult)
        assert result.content_type == expected_type
        assert result.backend_used == expected_backend
        if expected_backend == "error":
            assert result.message == ERROR_MESSAGE

    def test_docx_returns_real_html(
        self, office_file, all_backends_down
    ) -> None:
        """docx 降级到 mammoth HTML：content 为包含正文文本的 HTML 字符串。"""
        with all_backends_down:
            result = OfficeConverter.convert(
                {"path": office_file("docx"), "suffix": "docx", "is_dir": False}
            )

        assert result.content_type == "html"
        assert result.backend_used == "pure-python"
        assert isinstance(result.content, str)
        assert "Hello 世界" in result.content

    def test_pptx_returns_page_outline(
        self, office_file, all_backends_down
    ) -> None:
        """pptx 降级为逐页大纲：含页标记与各页文本。"""
        with all_backends_down:
            result = OfficeConverter.convert(
                {"path": office_file("pptx"), "suffix": "pptx", "is_dir": False}
            )

        assert result.content_type == "outline"
        assert result.backend_used == "pure-python"
        assert "--- 第 1 页 ---" in result.content
        assert "--- 第 2 页 ---" in result.content
        assert "Slide One" in result.content
        assert "Slide Two" in result.content

    def test_xlsx_returns_tsv_table(
        self, office_file, all_backends_down
    ) -> None:
        """xlsx 降级为 TSV 表格：行以换行、单元格以制表符分隔。"""
        with all_backends_down:
            result = OfficeConverter.convert(
                {"path": office_file("xlsx"), "suffix": "xlsx", "is_dir": False}
            )

        assert result.content_type == "table"
        assert result.backend_used == "pure-python"
        assert result.truncated is False
        rows = result.content.split("\n")
        assert rows[0] == "名称\t数量"
        assert rows[1] == "苹果\t3"
        assert rows[2] == "香蕉\t5"

    def test_legacy_defensive_error_from_pure_python_direct(
        self, office_file
    ) -> None:
        """防御性兜底：直接调用 ``_convert_pure_python`` 时 legacy 返回安装提示。"""
        result = OfficeConverter._convert_pure_python(
            {"path": office_file("doc"), "suffix": "doc", "is_dir": False},
            "doc",
        )

        assert result.content_type == "error"
        assert result.backend_used == "pure-python"
        assert result.message == ERROR_MESSAGE


# ===========================================================================
# Metis E7：xlsx 行 / 列上限与「已截断」标记
# ===========================================================================


class TestXlsxTruncationLimits:
    """超过 5000 行 / 200 列时应截断并设置 ``truncated=True`` + 「已截断」。"""

    def test_more_than_5000_rows_truncated(
        self, large_xlsx_path, all_backends_down
    ) -> None:
        """5200 行 → 返回行数 ≤5000 且标记截断。"""
        with all_backends_down:
            result = OfficeConverter.convert(
                {"path": large_xlsx_path, "suffix": "xlsx", "is_dir": False}
            )

        assert result.content_type == "table"
        assert result.truncated is True
        assert "已截断" in result.message
        row_count = result.content.count("\n") + 1
        assert row_count <= OfficeConverter._XLSX_MAX_ROWS
        assert row_count == 5000
        assert "row0" in result.content
        assert "row5199" not in result.content

    def test_more_than_200_columns_truncated(
        self, wide_xlsx_path, all_backends_down
    ) -> None:
        """250 列 → 每行仅保留前 200 列且标记截断。"""
        with all_backends_down:
            result = OfficeConverter.convert(
                {"path": wide_xlsx_path, "suffix": "xlsx", "is_dir": False}
            )

        assert result.content_type == "table"
        assert result.truncated is True
        assert "已截断" in result.message
        first_row = result.content.split("\n", 1)[0]
        assert first_row.count("\t") + 1 <= OfficeConverter._XLSX_MAX_COLS
        assert "col0" in first_row
        assert "col249" not in first_row


# ===========================================================================
# Metis B4-4：单一库缺失时仅对应后缀失败，其余后缀不受影响
# ===========================================================================


class TestMissingLibraryIsolation:
    """库缺失（ImportError）时对应后缀返回错误结果，其他后缀仍正常工作。"""

    @staticmethod
    def _block_import(monkeypatch, blocked: str) -> None:
        """屏蔽指定顶层模块的导入，并从 sys.modules 清除既有缓存。"""
        real_import = builtins.__import__

        def _fake_import(name, *args, **kwargs):
            if name == blocked or name.startswith(blocked + "."):
                raise ImportError(f"{blocked} blocked for test")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", _fake_import)
        monkeypatch.delitem(sys.modules, blocked, raising=False)

    def test_only_mammoth_blocked_docx_fails_xlsx_ok(
        self, monkeypatch, office_file, all_backends_down
    ) -> None:
        """仅 mammoth 缺失：docx 优雅失败，xlsx 仍成功。"""
        docx_file = office_file("docx")
        xlsx_file = office_file("xlsx")
        self._block_import(monkeypatch, "mammoth")
        with all_backends_down:
            docx_result = OfficeConverter.convert(
                {"path": docx_file, "suffix": "docx", "is_dir": False}
            )
            xlsx_result = OfficeConverter.convert(
                {"path": xlsx_file, "suffix": "xlsx", "is_dir": False}
            )

        assert docx_result.content_type == "error"
        assert docx_result.backend_used == "error"
        assert "mammoth" in docx_result.message
        assert xlsx_result.content_type == "table"
        assert xlsx_result.backend_used == "pure-python"

    def test_only_openpyxl_blocked_xlsx_fails_docx_ok(
        self, monkeypatch, office_file, all_backends_down
    ) -> None:
        """仅 openpyxl 缺失：xlsx 优雅失败，docx 仍成功。"""
        xlsx_file = office_file("xlsx")
        docx_file = office_file("docx")
        self._block_import(monkeypatch, "openpyxl")
        with all_backends_down:
            xlsx_result = OfficeConverter.convert(
                {"path": xlsx_file, "suffix": "xlsx", "is_dir": False}
            )
            docx_result = OfficeConverter.convert(
                {"path": docx_file, "suffix": "docx", "is_dir": False}
            )

        assert xlsx_result.content_type == "error"
        assert xlsx_result.backend_used == "error"
        assert "openpyxl" in xlsx_result.message
        assert docx_result.content_type == "html"
        assert docx_result.backend_used == "pure-python"


# ===========================================================================
# 文件缺失 / 损坏 → 返回结果，绝不抛出
# ===========================================================================


class TestMissingAndCorruptFiles:
    """缺失或损坏的文件必须返回错误 / 降级结果，绝不向外抛出。"""

    def test_missing_docx_file_returns_degraded_result(
        self, all_backends_down
    ) -> None:
        """不存在的 docx 路径返回降级结果（保留 pure-python 后端），不抛出。"""
        with all_backends_down:
            result = OfficeConverter.convert(
                {
                    "path": "C:/definitely/not/here/missing.docx",
                    "suffix": "docx",
                    "is_dir": False,
                }
            )

        assert result.backend_used == "pure-python"
        assert result.content_type == "html"
        assert result.content == ""
        assert "文件不存在" in result.message

    def test_corrupt_docx_returns_error_result(self, tmp_path, all_backends_down):
        """损坏的 docx（非 zip 内容）→ 错误结果，不抛出。"""
        corrupt = tmp_path / "corrupt.docx"
        corrupt.write_bytes(b"this is definitely not a docx zip file")
        with all_backends_down:
            result = OfficeConverter.convert(
                {"path": str(corrupt), "suffix": "docx", "is_dir": False}
            )

        assert result.content_type == "error"
        assert result.backend_used == "error"
        assert result.message

    def test_corrupt_xlsx_returns_error_result(self, tmp_path, all_backends_down):
        """损坏的 xlsx（垃圾字节）→ 错误结果，不抛出。"""
        corrupt = tmp_path / "corrupt.xlsx"
        corrupt.write_bytes(b"garbage bytes, not an xlsx archive")
        with all_backends_down:
            result = OfficeConverter.convert(
                {"path": str(corrupt), "suffix": "xlsx", "is_dir": False}
            )

        assert result.content_type == "error"
        assert result.backend_used == "error"
        assert result.message

    def test_missing_path_key_does_not_raise(self, all_backends_down) -> None:
        """缺少 ``path`` 键时不抛出，返回降级结果。"""
        with all_backends_down:
            result = OfficeConverter.convert({"suffix": "pptx"})
        assert result.backend_used == "pure-python"
        assert result.content_type == "outline"
        assert result.content == ""


# ===========================================================================
# 惰性导入：模块加载不得引入任何 Office 库（Metis B4-4）
# ===========================================================================


class TestLazyImports:
    """导入 ``office_converter`` 模块不得连带加载 mammoth / openpyxl / pptx。"""

    def test_module_import_does_not_load_office_libs(self, monkeypatch) -> None:
        """重新导入模块后三个库均不在 ``sys.modules`` 中。"""
        module_name = "freeassetfilter.services.office_converter"
        for name in (module_name, *_LAZY_LIBS):
            monkeypatch.delitem(sys.modules, name, raising=False)

        module = importlib.import_module(module_name)

        for name in _LAZY_LIBS:
            assert name not in sys.modules, (
                f"导入 {module_name} 不应连带导入 {name}"
            )
        assert hasattr(module, "OfficeConverter")

    def test_lazy_import_happens_only_inside_method(
        self, office_file, all_backends_down
    ) -> None:
        """真实转换前不加载 mammoth；转换后才出现在 ``sys.modules``。"""
        for name in _LAZY_LIBS:
            sys.modules.pop(name, None)

        with all_backends_down:
            # 先确认调用前 mammoth 不在 sys.modules（未因模块导入被加载）。
            assert "mammoth" not in sys.modules
            result = OfficeConverter.convert(
                {"path": office_file("docx"), "suffix": "docx", "is_dir": False}
            )

        assert result.content_type == "html"
        assert "mammoth" in sys.modules
