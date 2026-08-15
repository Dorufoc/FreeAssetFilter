#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
FreeAssetFilter v1.0

Copyright (c) 2026 Dorufoc <dorufoc@outlook.com>

协议说明：本软件基于 AGPL-3.0 协议开源
1. 个人非商业使用：需保留本注释及开发者署名；

项目地址：https://github.com/Dorufoc/FreeAssetFilter
许可协议：https://github.com/Dorufoc/FreeAssetFilter/blob/main/LICENSE

Office 转换分派服务
普通类（非单例），classmethod 风格（镜像 ImageDecoderService）。
能力探测（LibreOffice / MS Office / WPS COM）+ 后端分派：
按「LibreOffice → COM → 纯 Python」顺序自动降级，每个后缀定义允许的后端集合。
"""

from __future__ import annotations

import shutil
import subprocess
import tempfile
import threading
import time
import uuid
from collections.abc import Mapping
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from types import MappingProxyType

from freeassetfilter.utils.app_logger import debug, warning

# ── 结果对象 ─────────────────────────────────────────────────────────────


@dataclass(frozen=True)
class ConversionResult:
    """一次 Office 转换 / 降级尝试的不可变结果。

    Attributes
    ----------
    content_type :
        结果内容类型：``"pdf"``（LO/COM 转 PDF）、``"html"``（docx 降级）、
        ``"outline"``（pptx 降级）、``"table"``（xlsx 降级）、``"error"``。
    content :
        结果内容：转换产物路径（``Path``）或内联内容字符串（``str``）。
    backend_used :
        实际生效的后端标识：``"libreoffice"`` / ``"com"`` / ``"pure-python"``
        / ``"error"``。
    truncated :
        内容是否因行 / 列上限被截断（T7 xlsx 表格降级使用）。
    message :
        附加说明或错误提示文案。
    """

    content_type: str
    content: str | Path
    backend_used: str
    truncated: bool = False
    message: str = ""


# ── 常量（不可变） ───────────────────────────────────────────────────────
# 唯一例外：T5/T9 LibreOffice 后端持有模块级可变状态（_LO_CONVERSION_LOCK
# 转换串行锁与 _ACTIVE_LO_POPEN Popen 句柄注册表，见本文件下方）。


# 本服务支持的全部 Office 后缀（不含前导点，小写）。
SUPPORTED_SUFFIXES: frozenset[str] = frozenset({
    "docx", "pptx", "xlsx", "doc", "xls", "ppt",
})

# OOXML 现代格式：无任何外部后端时仍可纯 Python 降级。
MODERN_SUFFIXES: frozenset[str] = frozenset({"docx", "pptx", "xlsx"})

# OOXML 之前的旧二进制格式：无 LO/COM 时仅提示安装，不提供纯 Python 降级。
LEGACY_SUFFIXES: frozenset[str] = frozenset({"doc", "xls", "ppt"})

# 每个后缀允许的后端集合（元组顺序即分派优先级）。只读常量，运行期不修改。
_ALLOWED_BACKENDS: Mapping[str, tuple[str, ...]] = MappingProxyType({
    "docx": ("libreoffice", "com", "pure-python"),
    "pptx": ("libreoffice", "com", "pure-python"),
    "xlsx": ("libreoffice", "com", "pure-python"),
    "doc": ("libreoffice", "com"),
    "xls": ("libreoffice", "com"),
    "ppt": ("libreoffice", "com"),
})

# 纯 Python 降级（T7）各后缀对应的内容类型。只读常量。
_PURE_PYTHON_CONTENT_TYPE: Mapping[str, str] = MappingProxyType({
    "docx": "html",
    "pptx": "outline",
    "xlsx": "table",
})

# COM 探测的 ProgID 顺序（镜像 ``tests/conftest.py`` 的 ``com_available`` fixture）。
_COM_PROG_IDS: tuple[str, ...] = (
    "Word.Application",
    "Excel.Application",
    "PowerPoint.Application",
    "Kwps.Application",
    "Ket.Application",
    "Kwpp.Application",
)

# legacy 格式（doc/xls/ppt）在 LO/COM 均不可用时的提示文案（精确匹配计划措辞）。
ERROR_MESSAGE = "请安装 LibreOffice 或 Microsoft Office/WPS 以获得完整预览"

# LibreOffice 后端默认转换超时（秒）。T5/T9：测试通过注入更短的 ``timeout``
# 验证超时路径；soffice 卡死时由 ``_convert_with_libreoffice`` 终止并强杀。
SOFFICE_CONVERSION_TIMEOUT: float = 30.0

# LibreOffice 转换串行化锁：soffice 为单实例程序，同一时刻只允许一个转换
# 运行（跨线程）。模块级锁包裹「启动 + 等待」全程（T5 串行化测试契约）。
_LO_CONVERSION_LOCK: threading.Lock = threading.Lock()

# T9 取消 seam：正在进行的 soffice ``Popen`` 句柄注册表，键为执行线程的
# ``threading.get_ident()``。转换完成 / 失败后必须在 ``finally`` 中清空对应
# 键，避免泄漏句柄（``get_active_lo_popen()`` / ``clear_active_lo_popen()``）。
_ACTIVE_LO_POPEN: dict[int, subprocess.Popen] = {}


def get_active_lo_popen(thread_id: int | None = None) -> subprocess.Popen | None:
    """
    读取当前（或指定）线程正在进行的 soffice 子进程句柄（T9 取消 seam）。

    Parameters
    ----------
    thread_id : int | None
        线程标识（``threading.get_ident()``）；``None`` 时读取当前线程。

    Returns
    -------
    subprocess.Popen | None
        正在进行的 soffice ``Popen`` 句柄；无则返回 ``None``。
    """
    key = threading.get_ident() if thread_id is None else thread_id
    return _ACTIVE_LO_POPEN.get(key)


def clear_active_lo_popen(thread_id: int | None = None) -> None:
    """
    清除当前（或指定）线程的 soffice 子进程句柄（T9 取消 seam）。

    Parameters
    ----------
    thread_id : int | None
        线程标识（``threading.get_ident()``）；``None`` 时清除当前线程。
    """
    key = threading.get_ident() if thread_id is None else thread_id
    _ACTIVE_LO_POPEN.pop(key, None)


def _resolve_soffice_binary() -> Path | None:
    """
    在候选目录中定位 soffice 可执行文件。绝不抛出。

    复用 T2 ``soffice_paths()`` 返回的「目录」候选列表，逐个探测
    ``soffice.exe`` / ``soffice.com``；惰性导入为保证测试 monkeypatch
    ``core._paths.soffice_paths`` 生效（探测异常一律视为不可用）。

    Returns
    -------
    Path | None
        找到的 soffice 可执行文件路径；未找到或探测异常时返回 ``None``。
    """
    try:
        from freeassetfilter.core._paths import soffice_paths
        candidates = soffice_paths()
    except Exception as e:
        warning(f"[OfficeConverter] 解析 soffice 路径失败: {e}")
        return None
    for base in candidates:
        base_dir = Path(base)
        if not base_dir.is_dir():
            continue
        for name in ("soffice.exe", "soffice.com"):
            candidate = base_dir / name
            if candidate.is_file():
                return candidate
    return None


def _lo_degraded_result(message: str) -> ConversionResult:
    """
    源文件缺失 / 路径无效时的 LibreOffice 降级结果（T4 分派门禁契约）。

    保持 ``content_type="pdf"`` 与 ``backend_used="libreoffice"``（镜像 T7
    ``_degraded_result`` 约定），实际内容为空，失败原因写入 *message*。

    Parameters
    ----------
    message : str
        降级原因提示文案。

    Returns
    -------
    ConversionResult
        ``content=""`` 且携带提示消息的降级结果。
    """
    return ConversionResult(
        content_type="pdf",
        content="",
        backend_used="libreoffice",
        message=message,
    )


def _lo_error_result(message: str) -> ConversionResult:
    """
    生成 LibreOffice 后端的错误结果（``content_type="error"``）。

    Parameters
    ----------
    message : str
        明确的中文错误提示文案。

    Returns
    -------
    ConversionResult
        ``backend_used="error"`` 的错误结果。
    """
    return ConversionResult(
        content_type="error",
        content="",
        backend_used="error",
        message=message,
    )


def _terminate_lo_process(proc: subprocess.Popen) -> None:
    """
    超时后确保 soffice 子进程死亡（terminate → 宽限 → kill）。

    ``terminate()`` 后等待 1 秒宽限期；仍存活则 ``kill()`` 强杀并等待
    子进程退出，保证退出后 ``proc.poll()`` 非 ``None``。

    Parameters
    ----------
    proc : subprocess.Popen
        超时的 soffice 子进程句柄。
    """
    proc.terminate()
    try:
        proc.wait(timeout=1.0)
    except subprocess.TimeoutExpired:
        proc.kill()
        proc.wait()


def _remove_temp_profile(
    profile_dir: Path,
    attempts: int = 3,
    delay: float = 0.5,
) -> None:
    """
    尽力删除临时 LibreOffice profile 目录（容忍 Windows 文件锁）。

    最多重试 *attempts* 次、每次间隔 *delay* 秒；删除彻底失败仅记录告警
    日志，绝不抛出、绝不阻断转换流程（残留目录可容忍）。

    Parameters
    ----------
    profile_dir : Path
        待删除的临时 profile 目录。
    attempts : int
        删除重试次数上限（默认 3）。
    delay : float
        两次重试之间的间隔秒数（默认 0.5）。
    """
    for attempt in range(attempts):
        try:
            shutil.rmtree(profile_dir)
            return
        except Exception as e:
            if attempt + 1 < attempts:
                time.sleep(delay)
            else:
                warning(
                    f"[OfficeConverter] 清理临时 LibreOffice profile 失败: {e}"
                )


class OfficeConverter:
    """
    Office 转换分派服务（普通类，非单例，classmethod 风格）。

    职责：
    1. 能力探测 —— 每次 ``convert()`` 调用时运行时重新评估 LibreOffice /
       MS Office / WPS COM 的可用性，探测函数永不抛出；
    2. 后端分派 —— 按「LO → COM → 纯 Python」顺序为每个后缀选择允许的
       后端；doc/xls/ppt（legacy）在 LO/COM 均不可用时返回安装提示后端。

    本类无 ``__init__`` 实例状态，全部入口为 classmethod；缓存清理
    （T8）通过 ``_maybe_cleanup_cache()`` 在 ``convert()`` 入口幂等触发。

    使用方式::

        result = OfficeConverter.convert({"path": "a.docx", "suffix": "docx"})
        if result.backend_used == "error":
            print(result.message)
    """

    # ── 公共 API ─────────────────────────────────────────────────────

    @classmethod
    def convert(cls, file_info: dict) -> ConversionResult:
        """
        按「缓存命中 → LO → COM → 纯 Python」顺序处理 *file_info*。

        先查缓存（T8 缓存 + 本任务接线）：命中直接返回
        ``backend_used="cache"`` 的 PDF 结果，**不调用任何转换后端**；
        未命中则按 LO → COM → 纯 Python 顺序分派，并在产物为 PDF 时
        写入缓存供下次复用。

        Parameters
        ----------
        file_info : dict
            与 ``PreviewerRegistry.get_previewer_class`` 契约一致：至少含
            ``"suffix"``（文件扩展名，可带前导点或大小写混用）与 ``"path"``
            （T4 分派暂不读取，T5/T6/T7 后端使用）。

        Returns
        -------
        ConversionResult
            永不抛出异常。无任何可用后端时返回错误提示结果
            （``backend_used == "error"``）。

        Notes
        -----
        缓存清理（T8）在方法入口通过 ``_maybe_cleanup_cache()`` 幂等触发，
        同时惰性启动周期清理线程；本服务为 classmethod 风格、无
        ``__init__``，因此清理与缓存读写绝不依赖构造器。
        """
        cls._maybe_cleanup_cache()

        # 缓存命中短路（不调用任何后端）
        cached = cls._get_cache_path_safe(file_info)
        if cached is not None:
            return ConversionResult(
                content_type="pdf",
                content=cached,
                backend_used="cache",
                message="命中 Office 转换缓存",
            )

        suffix = cls._normalize_suffix(file_info)
        if suffix not in _ALLOWED_BACKENDS:
            return cls._unsupported_result(suffix)

        for backend in _ALLOWED_BACKENDS[suffix]:
            result = cls._try_backend(backend, file_info, suffix)
            if result is not None:
                # 仅 PDF 产物（LO/COM）落缓存，pure-python 文本产物不缓存
                return cls._cache_pdf_result(file_info, result)

        # legacy 格式（doc/xls/ppt）在 LO/COM 均不可用时走到这里。
        return cls._error_backend(suffix)

    # ── 能力探测（永不抛出；探测放服务内供运行时重新评估） ─────────────

    @classmethod
    def _soffice_available(cls) -> bool:
        """
        探测 LibreOffice soffice 是否可用。永不抛出。

        复用 T2 的 ``freeassetfilter.core._paths.soffice_paths()``：
        任一候选目录存在 ``soffice.exe`` / ``soffice.com`` 即视为可用；
        导入或调用失败一律视为不可用（返回 ``False``）。

        Returns
        -------
        bool
            ``True`` 表示存在可达的 soffice 可执行文件。
        """
        try:
            from freeassetfilter.core._paths import soffice_paths
        except (ImportError, AttributeError) as e:
            debug(f"[OfficeConverter] 无法导入 soffice_paths: {e}")
            return False

        try:
            candidates = soffice_paths()
        except Exception as e:
            debug(f"[OfficeConverter] soffice_paths() 探测异常: {e}")
            return False

        return any(
            Path(p).is_dir()
            and any(
                (Path(p) / name).is_file()
                for name in ("soffice.exe", "soffice.com")
            )
            for p in candidates
        )

    @classmethod
    def _com_available(cls) -> bool:
        """
        探测 MS Office / WPS COM 组件是否可用。永不抛出。

        镜像 ``tests/conftest.py`` 的 ``com_available`` fixture：优先注册表
        探测 6 个 ProgID（不启动任何进程），仅当注册表探测不可行时回退到
        ``win32com.client.Dispatch`` try/except（finally 中尽力 ``Quit()``）。

        Returns
        -------
        bool
            任一 Office/WPS ProgID 可实例化则为 ``True``，否则为 ``False``。
        """
        prog_ids = _COM_PROG_IDS

        # 优先注册表探测
        try:
            import winreg
        except ImportError:
            winreg = None

        if winreg is not None:
            try:
                for prog_id in prog_ids:
                    try:
                        winreg.OpenKey(winreg.HKEY_CLASSES_ROOT, prog_id)
                        return True
                    except OSError:
                        continue
            except Exception as e:
                debug(f"[OfficeConverter] 注册表探测 COM 异常: {e}")
                return False

        # 回退：COM Dispatch 探测（可能短暂启动隐藏实例）
        try:
            import win32com.client
        except Exception:
            return False

        app = None
        for prog_id in prog_ids:
            try:
                app = win32com.client.Dispatch(prog_id)
                return True
            except Exception:
                continue
            finally:
                if app is not None:
                    try:
                        app.Quit()
                    except Exception:
                        pass
                    app = None

        return False

    # ── 分派内部实现 ─────────────────────────────────────────────────

    @classmethod
    def _try_backend(
        cls,
        backend: str,
        file_info: dict,
        suffix: str,
    ) -> ConversionResult | None:
        """
        尝试单个 *backend*：先探测可用性，可用则调用对应转换方法。

        探测函数抛出的任何异常都在此捕获并视为不可用（返回 ``None``，
        由调用方降级到下一级），绝不向上传播。纯 Python 后端在 T4 阶段
        恒可用（占位实现，T7 替换为真实提取并在依赖缺失时返回 ``None``）。

        Parameters
        ----------
        backend : str
            ``"libreoffice"`` / ``"com"`` / ``"pure-python"``。
        file_info : dict
            与 ``convert()`` 一致的文件信息。
        suffix : str
            归一化小写后缀（无前导点）。

        Returns
        -------
        ConversionResult | None
            后端不可用或探测失败时返回 ``None``。
        """
        if backend == "libreoffice":
            try:
                available = cls._soffice_available()
            except Exception as e:
                warning(f"[OfficeConverter] LibreOffice 探测异常，视为不可用: {e}")
                available = False
            if not available:
                return None
            return cls._convert_with_libreoffice(file_info, suffix)

        if backend == "com":
            try:
                available = cls._com_available()
            except Exception as e:
                warning(f"[OfficeConverter] COM 探测异常，视为不可用: {e}")
                available = False
            if not available:
                return None
            return cls._convert_with_com(file_info, suffix)

        # pure-python
        return cls._convert_pure_python(file_info, suffix)

    # ── 各后端转换方法（T4 占位；真实实现由 T5/T6/T7 替换） ────────────

    @classmethod
    def _convert_with_libreoffice(
        cls,
        file_info: dict,
        suffix: str = "",
        *,
        timeout: float = SOFFICE_CONVERSION_TIMEOUT,
    ) -> ConversionResult:
        """
        LibreOffice 后端：通过 soffice 无头子进程把 Office 文档转为 PDF。

        执行流程（Metis B3/D5/D8）：
        1. 输入校验 —— 路径缺失或源文件不存在时返回保留 ``"pdf"/"libreoffice"``
           标识的降级结果（镜像 T7 ``_degraded_result`` 约定，兼容 T4 分派
           路由测试），绝不启动子进程；
        2. 二进制解析 —— ``_resolve_soffice_binary()`` 复用 T2
           ``soffice_paths()`` 候选目录探测 ``soffice.exe``/``soffice.com``，
           未找到即错误结果（不抛出）；
        3. 串行化 —— 模块级 ``_LO_CONVERSION_LOCK`` 包裹「启动 + 等待」全程
           （soffice 为单实例程序，同一时刻只允许一个转换运行）；
        4. 临时目录 —— ``tempfile.mkdtemp`` 创建 ``faf_lo_*`` profile 与
           ``faf_lo_out_*`` 输出目录；profile 在 ``finally`` 中尽力删除
           （3 次 × 0.5s 重试，容忍 Windows 文件锁，失败仅告警）；
        5. 取消 seam —— 子进程句柄注册到模块级
           ``_ACTIVE_LO_POPEN[当前线程 ident]``（T9 集成点），无论成败都在
           ``finally`` 中 ``clear_active_lo_popen``；
        6. 超时 —— ``proc.wait(timeout=timeout)`` 超时后 ``terminate()`` →
           短暂宽限 → ``kill()`` → ``wait()``，确保子进程已死。

        Parameters
        ----------
        file_info : dict
            与 ``convert()`` 一致的文件信息（至少含 ``"path"``）。
        suffix : str
            归一化小写后缀（无前导点），仅作兼容参数，不影响输出命名
            （输出名始终为源文件 basename + ``.pdf``）。
        timeout : float
            转换超时秒数，默认 ``SOFFICE_CONVERSION_TIMEOUT``（30 秒）。

        Returns
        -------
        ConversionResult
            成功时 ``content`` 为产物 PDF 的 ``Path``，``content_type="pdf"``，
            ``backend_used="libreoffice"``；源文件缺失返回保留后端标识的降级
            结果；其余失败（二进制缺失 / 超时 / 非零退出 / 未产出 PDF /
            启动失败）返回 ``content_type="error"`` 的错误结果。所有路径
            绝不抛出。
        """
        path = cls._extract_path(file_info)
        if path is None:
            return _lo_degraded_result("LibreOffice 转换失败：无法获取文件路径")
        if not Path(path).is_file():
            return _lo_degraded_result(
                f"LibreOffice 转换失败：源文件不存在：{path}"
            )

        binary = _resolve_soffice_binary()
        if binary is None:
            return _lo_error_result(
                "LibreOffice 转换失败：未检测到 soffice（请安装 LibreOffice）"
            )

        # soffice 为单实例程序：模块级锁包裹「启动 + 等待」全程，保证并发
        # 转换串行执行（T5 测试用阻塞假 soffice 断言启动时间差 >= 阻塞时长）。
        with _LO_CONVERSION_LOCK:
            thread_id = threading.get_ident()
            profile_dir: Path | None = None
            out_dir: Path | None = None
            try:
                profile_dir = Path(tempfile.mkdtemp(prefix="faf_lo_"))
                out_dir = Path(tempfile.mkdtemp(prefix="faf_lo_out_"))
                output_pdf = out_dir / f"{Path(path).stem}.pdf"
                profile_flag = (
                    "-env:UserInstallation=file:///"
                    + profile_dir.as_posix().lstrip("/")
                )
                command = [
                    str(binary),
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(out_dir),
                    profile_flag,
                    str(path),
                ]

                try:
                    proc = subprocess.Popen(
                        command,
                        stdout=subprocess.DEVNULL,
                        stderr=subprocess.DEVNULL,
                    )
                except Exception as e:
                    warning(f"[OfficeConverter] LibreOffice 启动失败: {e}")
                    return _lo_error_result(f"LibreOffice 转换失败：{e}")

                # T9 取消 seam：句柄注册到当前线程，成功/失败都必须清理。
                _ACTIVE_LO_POPEN[thread_id] = proc
                try:
                    try:
                        proc.wait(timeout=timeout)
                    except subprocess.TimeoutExpired:
                        _terminate_lo_process(proc)
                        return _lo_error_result(
                            f"LibreOffice 转换超时（超过 {timeout} 秒），"
                            "已终止 soffice 进程"
                        )
                    if proc.returncode != 0:
                        return _lo_error_result(
                            "LibreOffice 转换失败：soffice 退出码 "
                            f"{proc.returncode}"
                        )
                    if not output_pdf.is_file() or output_pdf.stat().st_size == 0:
                        return _lo_error_result(
                            "LibreOffice 转换失败：soffice 退出但未生成 PDF 文件"
                        )
                    return ConversionResult(
                        content_type="pdf",
                        content=output_pdf,
                        backend_used="libreoffice",
                    )
                finally:
                    clear_active_lo_popen(thread_id)
            finally:
                if profile_dir is not None:
                    _remove_temp_profile(profile_dir)

    @classmethod
    def _convert_with_com(
        cls,
        file_info: dict,
        suffix: str,
    ) -> ConversionResult:
        """
        COM 后端（T6）：通过 MS Office / WPS COM 自动化将文档导出为 PDF。

        Metis A2（非协商项）—— 每次任务新建一个专用线程执行 COM 工作：
        ``pythoncom.CoInitialize()`` 是该任务线程的第一条语句，
        ``pythoncom.CoUninitialize()`` 放入 ``finally``；全部 COM 对象在
        同一任务线程内创建 / 使用 / 释放，绝不跨线程传递。Metis E8 用
        ``threading.get_ident()`` sentinel 断言 CoInitialize 运行在任务线程
        而非调用线程，且失败路径也要执行 CoUninitialize。

        Metis C4（仅顺序，无 WPS 特判）—— ProgID 顺序：doc/docx 依次
        ``Word.Application`` → ``Kwps.Application``；xls/xlsx 依次
        ``Excel.Application`` → ``Ket.Application``；ppt/pptx 依次
        ``PowerPoint.Application`` → ``Kwpp.Application``。任一可实例化即
        使用，不做品牌区分。

        pywin32（pythoncom / win32com.client）惰性导入只在方法内部；
        导入失败返回错误结果。自动化设置：``DisplayAlerts=False``（真实 QA
        验证 PowerPoint 接受）、``Visible=False``（仅 Word/Excel 允许，PowerPoint
        禁止隐藏主窗口故改用 ``WindowState=2`` 最小化）、``ReadOnly=True``；
        ``app.Quit()`` 收尾在 ``finally``。PowerPoint 换用 ``SaveAs(路径, 32)``
        （真实 QA 验证 win32com 动态分派不接受 ``PrintRange`` 参数）。

        Parameters
        ----------
        file_info : dict
            与 ``convert()`` 一致的文件信息（至少含 ``"path"``）。
        suffix : str
            归一化小写后缀（无前导点），决定 ProgID 顺序。

        Returns
        -------
        ConversionResult
            成功时 ``content`` 为产物 PDF 的 ``Path``，``content_type="pdf"``；
            源文件缺失时返回保留 ``"com"``/``"pdf"`` 标识的降级结果（镜像
            ``_degraded_result`` 约定，兼容 T4 分派路由测试）；
            其余失败返回 ``content_type="error"`` 的错误结果，绝不抛出。
        """
        path = cls._extract_path(file_info)
        if path is None:
            return cls._com_error_result("COM 转换失败：无法获取文件路径")

        # 源文件不存在（如 T4 分派路由测试的占位路径）→ 返回保留后端标识的
        # 降级结果（镜像 ``_degraded_result`` 约定），不启动任务线程。
        if not Path(path).is_file():
            return ConversionResult(
                content_type="pdf",
                content="",
                backend_used="com",
                message="COM 转换失败：源文件不存在",
            )

        try:
            import pythoncom  # noqa: PLC0415
            import win32com.client  # noqa: PLC0415
        except ImportError:
            return cls._com_error_result(
                "COM 转换失败：pywin32 未安装（pip install pywin32）"
            )

        # Metis A2/E8：结果经线程安全的容器回传；线程目标第一句 CoInitialize。
        results: list[ConversionResult] = []
        task_started_at = datetime.now()

        def _com_task() -> None:
            pythoncom.CoInitialize()
            try:
                results.append(
                    cls._com_convert_in_task(path, suffix, task_started_at)
                )
            except Exception as e:
                results.append(cls._com_error_result(f"COM 转换失败：{e}"))
            finally:
                pythoncom.CoUninitialize()

        thread = threading.Thread(
            target=_com_task,
            name=f"office-com-{suffix}",
            daemon=True,
        )
        thread.start()
        thread.join()

        if not results:
            return cls._com_error_result("COM 转换失败：任务线程未返回结果")
        return results[0]

    @staticmethod
    def _com_error_result(message: str) -> ConversionResult:
        """
        生成 COM 后端的错误结果（``content_type="error"``）。

        Parameters
        ----------
        message : str
            明确的中文错误提示文案。

        Returns
        -------
        ConversionResult
            ``backend_used="error"`` 的错误结果。
        """
        return ConversionResult(
            content_type="error",
            content="",
            backend_used="error",
            message=message,
        )

    @classmethod
    def _com_convert_in_task(
        cls,
        path: str,
        suffix: str,
        task_started_at: datetime,
    ) -> ConversionResult:
        """
        在任务线程内完成真实的 COM 自动化转换。仅在任务线程执行。

        该方法的调用方（``_convert_with_com``）已完成 ``CoInitialize`` 并在
        ``finally`` 中调用 ``CoUninitialize``；本方法关注单次文档转换本身。

        Parameters
        ----------
        path : str
            源文档的字符串路径。
        suffix : str
            归一化小写后缀（无前导点）。
        task_started_at : datetime
            本次转换任务开始时间，用于孤儿进程清理时间窗下界（Metis B3）。

        Returns
        -------
        ConversionResult
            成功返回 ``content`` 为产物 PDF 的 ``Path``；失败返回错误结果。
        """
        import win32com.client  # noqa: PLC0415

        prog_ids = cls._com_prog_ids(suffix)
        if not prog_ids:
            return cls._com_error_result(f"COM 转换失败：不支持的后缀 {suffix}")

        app = None
        document = None
        try:
            # Metis C4：按顺序尝试 ProgID，任一成功即使用。
            for prog_id in prog_ids:
                try:
                    app = win32com.client.Dispatch(prog_id)
                    break
                except Exception as e:
                    warning(f"[OfficeConverter] COM 启动 {prog_id} 失败: {e}")
            if app is None:
                return cls._com_error_result(
                    "COM 转换失败：未能启动 MS Office/WPS（请确认已安装）"
                )

            app.DisplayAlerts = False
            if suffix in ("ppt", "pptx"):
                # PowerPoint 禁止隐藏主窗口（COM 异常 -2147352567 ...
                # "Application.Visible : Invalid request. Hiding the
                # application window is not allowed."）。改用最小化窗口
                # 替代（ppWindowMinimized = 2），既不弹大窗口也不违反
                # PowerPoint 的限制（真实 QA 验证合法）。
                app.WindowState = 2
            else:
                app.Visible = False

            out_dir = Path(tempfile.gettempdir())
            out_dir.mkdir(parents=True, exist_ok=True)
            output = out_dir / f"{Path(path).stem}_com_{uuid.uuid4().hex}.pdf"

            if suffix in ("doc", "docx"):
                document = app.Documents.Open(
                    path,
                    ReadOnly=True,
                )
                # wdExportFormatPDF = 17
                document.ExportAsFixedFormat(str(output), 17)
            elif suffix in ("xls", "xlsx"):
                document = app.Workbooks.Open(
                    path,
                    ReadOnly=True,
                    UpdateLinks=0,
                )
                # xlTypePDF = 0
                document.ActiveSheet.ExportAsFixedFormat(0, str(output))
            else:  # ppt / pptx
                document = app.Presentations.Open(
                    path,
                    ReadOnly=True,
                    WithWindow=False,
                )
                # ppSaveAsPDF = 32。真实 QA 验证：win32com 动态分派对
                # PowerPoint 的 SaveAs 没有 PrintRange 参数（传它报
                # "unexpected keyword argument 'PrintRange'"），故不传。
                document.SaveAs(str(output), 32)

            if not output.is_file() or output.stat().st_size == 0:
                return cls._com_error_result("COM 转换失败：未生成 PDF 产物")

            return ConversionResult(
                content_type="pdf",
                content=output,
                backend_used="com",
            )
        except Exception as e:
            warning(f"[OfficeConverter] COM 转换异常: {e}")
            return cls._com_error_result(f"COM 转换失败：{e}")
        finally:
            if document is not None:
                try:
                    document.Close(False)
                except Exception:
                    pass
            if app is not None:
                try:
                    app.Quit()
                except Exception:
                    pass
            cls._cleanup_orphan_processes(task_started_at)

    @staticmethod
    def _com_prog_ids(suffix: str) -> tuple[str, ...]:
        """
        按后缀返回 ProgID 尝试顺序（Metis C4，仅顺序、无品牌特判）。

        Parameters
        ----------
        suffix : str
            归一化小写后缀（无前导点）。

        Returns
        -------
        tuple[str, ...]
            MS Office ProgID 在前的有序元组；未知后缀返回空元组。
        """
        if suffix in ("doc", "docx"):
            return ("Word.Application", "Kwps.Application")
        if suffix in ("xls", "xlsx"):
            return ("Excel.Application", "Ket.Application")
        if suffix in ("ppt", "pptx"):
            return ("PowerPoint.Application", "Kwpp.Application")
        return ()

    @staticmethod
    def _cleanup_orphan_processes(task_started_at: datetime) -> None:
        """
        Metis B3：尽力清理本服务可能残留的 Office/WPS 孤儿进程。

        仅结束 **本次转换启动的隐藏实例** —— 判定条件（须同时满足）：
        - 进程名在 WINWORD / EXCEL / POWERPNT 之列；
        - 进程启动时间晚于本次转换任务开始（``task_started_at``）；
        - 进程没有主窗口标题（COM 隐藏实例 ``Visible=False`` 无窗口；用户
          手动打开的文档必然有标题）。

        找不到 / 无法区分 / powershell 不可用时保持安全（不终止任何进程），
        绝不影响用户的其它 Office 文档。所有操作均为 best-effort，
        任何失败静默忽略、仅记调试日志。

        Parameters
        ----------
        task_started_at : datetime
            本次转换任务开始时间，作为进程启动时间窗下界。
        """
        try:
            cutoff = task_started_at.strftime("%Y-%m-%dT%H:%M:%S")
            stdout = subprocess.run(
                [
                    "powershell",
                    "-NoProfile",
                    "-Command",
                    (
                        "Get-Process WINWORD,EXCEL,POWERPNT "
                        "-ErrorAction SilentlyContinue | Where-Object {"
                        f"$_.StartTime -ge [datetime]'{cutoff}' -and "
                        "$_.MainWindowTitle -eq ''} | "
                        "Stop-Process -Force -ErrorAction SilentlyContinue"
                    ),
                ],
                capture_output=True,
                timeout=15,
                check=False,
            )
            if stdout.returncode != 0:
                debug(
                    f"[OfficeConverter] 孤儿进程清理命令退出码 {stdout.returncode}"
                )
        except FileNotFoundError:
            debug("[OfficeConverter] powershell 不可用，跳过孤儿进程清理")
        except Exception as e:
            debug(f"[OfficeConverter] 孤儿进程清理失败（安全跳过）: {e}")

    @classmethod
    def _convert_pure_python(
        cls,
        file_info: dict,
        suffix: str,
    ) -> ConversionResult:
        """
        纯 Python 后端：docx→HTML、pptx→大纲、xlsx→表格；失败时返回错误结果。

        仅当 LibreOffice 与 COM 均不可用时被调用（现代格式 docx/pptx/xlsx）。
        全部第三方库（mammoth / python-pptx / openpyxl）都在各方法内部按需
        惰性导入（Metis B4-4）—— 模块加载不依赖任何 Office 库。

        xlsx 表格表示约定（T10 表格视图消费协议）：
        ``content`` 为 TSV（制表符分隔）字符串 —— 行以 ``\\n`` 分隔、单元格以
        ``\\t`` 分隔；每个单元格值统一转为 ``str``（``None`` → 空串），单元格内
        原有制表符 / 换行符替换为空格以保证 TSV 结构完整。行上限 ``_XLSX_MAX_ROWS``
        （5000）、列上限 ``_XLSX_MAX_COLS``（200）；任一超限时 ``truncated=True``
        且 ``message`` 含「已截断」。T10 按 ``content.split("\\n")`` → 每行
        ``split("\\t")`` 即可还原为表格。

        Parameters
        ----------
        file_info : dict
            与 ``convert()`` 一致的文件信息（至少含 ``"path"``）。
        suffix : str
            归一化小写后缀（无前导点）。

        Returns
        -------
        ConversionResult
            - docx：``content_type="html"``，``content`` 为 mammoth 生成的 HTML 字符串；
            - pptx：``content_type="outline"``，``content`` 为逐页纯文本大纲；
            - xlsx：``content_type="table"``，``content`` 为 TSV 字符串；
            - legacy / 未知后缀（防御性兜底，正常分派不会到达）：错误结果；
            - 依赖缺失 / 解析失败：``backend_used="error"`` 的错误结果，绝不抛出；
            - 文件缺失 / 路径无效：保持 ``content_type`` 与 ``backend_used`` 的
              降级结果（``content=""`` + 提示消息），绝不抛出。
        """
        if suffix == "docx":
            return cls._pure_python_docx(file_info)
        if suffix == "pptx":
            return cls._pure_python_pptx(file_info)
        if suffix == "xlsx":
            return cls._pure_python_xlsx(file_info)
        # legacy（doc/xls/ppt）或未知后缀的防御性兜底 —— 正常分派不会走到这里。
        return ConversionResult(
            content_type="error",
            content="",
            backend_used="pure-python",
            message=ERROR_MESSAGE,
        )

    # ── 纯 Python 后端各格式实现（T7；所有第三方库导入均为惰性） ───────

    # xlsx 表格提取上限（T7；与 ``_pure_python_xlsx`` 的 TSV 表示约定一致）。
    _XLSX_MAX_ROWS: int = 5000
    _XLSX_MAX_COLS: int = 200

    @classmethod
    def _pure_python_docx(cls, file_info: dict) -> ConversionResult:
        """
        docx → HTML（mammoth）。库缺失 / 文件问题一律返回结果而不抛出。

        Parameters
        ----------
        file_info : dict
            含 ``"path"`` 的文件信息。

        Returns
        -------
        ConversionResult
            ``content_type="html"`` 且 ``content`` 为 mammoth 生成的 HTML 字符串；
            失败时为错误 / 降级结果。
        """
        path = cls._extract_path(file_info)
        if path is None:
            return cls._degraded_result("docx", "未提供有效的文件路径")
        if not Path(path).is_file():
            return cls._degraded_result("docx", f"文件不存在或不可读：{path}")

        try:
            import mammoth
        except ImportError:
            warning("[OfficeConverter] mammoth 不可用，docx 纯 Python 降级失败")
            return cls._pure_python_error_result(
                "无法预览 docx 文档：缺少 mammoth 库，请安装依赖后重试"
            )

        try:
            converted = mammoth.convert_to_html(path)
        except Exception as e:
            warning(f"[OfficeConverter] docx 转换失败: {e}")
            return cls._pure_python_error_result(f"docx 文档解析失败：{e}")

        return ConversionResult(
            content_type="html",
            content=converted.value,
            backend_used="pure-python",
        )

    @classmethod
    def _pure_python_pptx(cls, file_info: dict) -> ConversionResult:
        """
        pptx → 逐页纯文本大纲（python-pptx）。失败时返回结果而不抛出。

        ``content`` 为逐页大纲：每页以 ``--- 第 N 页 ---`` 标记开头，页内文本行
        按形状 / 段落顺序拼接，页与页之间空行分隔。

        Parameters
        ----------
        file_info : dict
            含 ``"path"`` 的文件信息。

        Returns
        -------
        ConversionResult
            ``content_type="outline"``；失败时为错误 / 降级结果。
        """
        path = cls._extract_path(file_info)
        if path is None:
            return cls._degraded_result("pptx", "未提供有效的文件路径")
        if not Path(path).is_file():
            return cls._degraded_result("pptx", f"文件不存在或不可读：{path}")

        try:
            from pptx import Presentation
        except ImportError:
            warning("[OfficeConverter] python-pptx 不可用，pptx 纯 Python 降级失败")
            return cls._pure_python_error_result(
                "无法预览 pptx 演示文稿：缺少 python-pptx 库，请安装依赖后重试"
            )

        try:
            presentation = Presentation(path)
        except Exception as e:
            warning(f"[OfficeConverter] pptx 解析失败: {e}")
            return cls._pure_python_error_result(f"pptx 文件解析失败：{e}")

        pages: list[str] = []
        for page_no, slide in enumerate(presentation.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        lines.append(text)
            pages.append(f"--- 第 {page_no} 页 ---\n" + "\n".join(lines))

        return ConversionResult(
            content_type="outline",
            content="\n\n".join(pages),
            backend_used="pure-python",
        )

    @classmethod
    def _pure_python_xlsx(cls, file_info: dict) -> ConversionResult:
        """
        xlsx → TSV 表格字符串（openpyxl 只读模式）。失败时返回结果而不抛出。

        表格表示约定见 ``_convert_pure_python`` 的 docstring：``content`` 为
        TSV 字符串，行 / 列上限分别为 ``_XLSX_MAX_ROWS``（5000）与
        ``_XLSX_MAX_COLS``（200），任一超限时 ``truncated=True`` 且
        ``message`` 含「已截断」。不做排序 / 筛选 / 公式 / 编辑（Metis C2）。

        Parameters
        ----------
        file_info : dict
            含 ``"path"`` 的文件信息。

        Returns
        -------
        ConversionResult
            ``content_type="table"``；失败时为错误 / 降级结果。
        """
        path = cls._extract_path(file_info)
        if path is None:
            return cls._degraded_result("xlsx", "未提供有效的文件路径")
        if not Path(path).is_file():
            return cls._degraded_result("xlsx", f"文件不存在或不可读：{path}")

        try:
            import openpyxl
        except ImportError:
            warning("[OfficeConverter] openpyxl 不可用，xlsx 纯 Python 降级失败")
            return cls._pure_python_error_result(
                "无法预览 xlsx 表格：缺少 openpyxl 库，请安装依赖后重试"
            )

        max_rows = cls._XLSX_MAX_ROWS
        max_cols = cls._XLSX_MAX_COLS
        truncated = False
        lines: list[str] = []
        workbook = None
        try:
            workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
        except Exception as e:
            warning(f"[OfficeConverter] xlsx 解析失败: {e}")
            return cls._pure_python_error_result(f"xlsx 文件解析失败：{e}")

        try:
            worksheet = workbook.active
            for row_index, row in enumerate(
                worksheet.iter_rows(values_only=True), start=1
            ):
                if row_index > max_rows:
                    truncated = True
                    break
                cells = row if row is not None else ()
                if len(cells) > max_cols:
                    truncated = True
                    cells = cells[:max_cols]
                lines.append("\t".join(cls._sanitize_cell(v) for v in cells))
        except Exception as e:
            warning(f"[OfficeConverter] xlsx 读取失败: {e}")
            return cls._pure_python_error_result(f"xlsx 表格读取失败：{e}")
        finally:
            if workbook is not None:
                workbook.close()

        message = (
            f"已截断：仅显示前 {max_rows} 行 / 前 {max_cols} 列"
            if truncated
            else ""
        )
        return ConversionResult(
            content_type="table",
            content="\n".join(lines),
            backend_used="pure-python",
            truncated=truncated,
            message=message,
        )

    @staticmethod
    def _extract_path(file_info: dict) -> str | None:
        """
        从 *file_info* 提取字符串路径；缺失 / 非路径 / 无法转换时返回 ``None``。

        Parameters
        ----------
        file_info : dict
            含 ``"path"`` 键的文件信息。

        Returns
        -------
        str | None
            字符串路径；``file_info`` 非字典、缺少 ``"path"`` 或无法转为字符串
            时返回 ``None``。
        """
        if not isinstance(file_info, dict):
            return None
        path = file_info.get("path")
        if path is None:
            return None
        try:
            return str(path)
        except Exception:
            return None

    @staticmethod
    def _degraded_result(suffix: str, message: str) -> ConversionResult:
        """
        文件缺失 / 路径无效时的降级结果（保持内容类型与后端标识）。

        保持 ``content_type`` 与 ``backend_used="pure-python"`` 是为了兼容 T4
        分派测试（对不存在的占位路径仍断言路由到纯 Python 后端）；实际内容为空，
        失败原因写入 *message*。

        Parameters
        ----------
        suffix : str
            归一化小写后缀，决定保留的 ``content_type``。
        message : str
            降级原因提示文案。

        Returns
        -------
        ConversionResult
            ``content=""`` 且携带提示消息的降级结果。
        """
        return ConversionResult(
            content_type=_PURE_PYTHON_CONTENT_TYPE.get(suffix, "outline"),
            content="",
            backend_used="pure-python",
            message=message,
        )

    @staticmethod
    def _pure_python_error_result(message: str) -> ConversionResult:
        """
        依赖缺失或解析失败时的错误结果（``backend_used="error"``）。

        Parameters
        ----------
        message : str
            明确的错误提示文案（中文）。

        Returns
        -------
        ConversionResult
            ``content_type="error"`` 的错误结果。
        """
        return ConversionResult(
            content_type="error",
            content="",
            backend_used="error",
            message=message,
        )

    @staticmethod
    def _sanitize_cell(value: object) -> str:
        """
        将单元格值转为字符串并清洗制表符 / 换行，保证 TSV 结构完整。

        Parameters
        ----------
        value : object
            openpyxl 单元格原始值（可为 ``None``）。

        Returns
        -------
        str
            清洗后的单元格字符串；``None`` 返回空串。
        """
        if value is None:
            return ""
        return str(value).replace("\t", " ").replace("\r", " ").replace("\n", " ")

    # ── 错误 / 提示结果 ───────────────────────────────────────────────

    @classmethod
    def _error_backend(cls, suffix: str) -> ConversionResult:
        """
        legacy 格式（doc/xls/ppt）在 LO/COM 均不可用时的安装提示后端。

        Parameters
        ----------
        suffix : str
            归一化小写后缀，仅作日志 / 调试标识。

        Returns
        -------
        ConversionResult
            ``content_type="error"`` 且 ``message`` 为安装提示文案。
        """
        return ConversionResult(
            content_type="error",
            content="",
            backend_used="error",
            message=ERROR_MESSAGE,
        )

    @classmethod
    def _unsupported_result(cls, suffix: str) -> ConversionResult:
        """
        未知后缀或缺少后缀信息时的错误结果（``convert()`` 不抛出）。

        Parameters
        ----------
        suffix : str
            归一化小写后缀，用于错误消息；为空表示缺失后缀信息。

        Returns
        -------
        ConversionResult
            ``content_type="error"`` 的错误结果。
        """
        message = (
            f"不支持的 Office 文件格式: {suffix}"
            if suffix
            else "缺少文件后缀信息，无法分派 Office 转换"
        )
        return ConversionResult(
            content_type="error",
            content="",
            backend_used="error",
            message=message,
        )

    # ── T8 接入点（缓存清理 seam） ────────────────────────────────────

    @classmethod
    def _maybe_cleanup_cache(cls) -> None:
        """
        缓存清理幂等入口（T8 实现 + 周期清理惰性启动）。

        设计约束：本服务为 classmethod 风格、无 ``__init__``，因此缓存清理
        与周期清理线程的启动都必须在 ``convert()`` 入口调用本方法幂等触发，
        绝不依赖构造器。惰性导入 ``office_cache.cleanup_cache`` 并调用
        （7 天过期 + 大小上限驱逐），同时惰性启动周期清理 daemon 线程
        （默认每 30 分钟自动清理一次）；任何异常都静默降级，绝不阻断
        转换。重复调用无副作用（启动是幂等的）。
        """
        try:
            from freeassetfilter.services.office_cache import (
                _ensure_periodic_cleanup_started,
                cleanup_cache,
            )
        except ImportError:
            # office_cache 缺失时跳过清理，不影响转换主流程。
            return
        try:
            # 首次 convert 时惰性启动周期清理线程（幂等，应用生命周期内只启动一次）。
            _ensure_periodic_cleanup_started()
            cleanup_cache()
        except Exception:
            # 缓存清理失败绝不阻断转换（降级为不清理）。
            pass

    # ── 缓存读写接入（T14 需求接线） ─────────────────────────────────

    @staticmethod
    def _get_cache_path_safe(file_info: dict) -> Path | None:
        """
        查询 *file_info* 的缓存命中路径；任何异常/缺失信息都降级为未命中。

        惰性导入 ``office_cache.get_cache_path``（镜像 ``_maybe_cleanup_cache``
        的容错风格）：导入失败、调用抛异常或目录不可写一律返回 ``None``，
        绝不因缓存问题阻断转换。

        Parameters
        ----------
        file_info : dict
            与 ``convert()`` 一致的文件信息（含 ``"path"``）。

        Returns
        -------
        Path | None
            缓存 PDF 路径（命中），或 ``None``（未命中 / 缓存不可用）。
        """
        try:
            from freeassetfilter.services.office_cache import get_cache_path
        except ImportError:
            return None
        try:
            return get_cache_path(file_info)
        except Exception:
            return None

    @classmethod
    def _cache_pdf_result(cls, file_info: dict, result: ConversionResult) -> ConversionResult:
        """
        把 PDF 产物写入缓存；兼容非缓存路径（原样返回）。

        仅当 *result* 是真实的 PDF 产物（``content_type=="pdf"`` 且
        ``content`` 为 ``Path`` 且 ``backend_used`` 为 ``"libreoffice"`` /
        ``"com"``）时才调用 ``put_cache``；pure-python 后端的
        html/outline/table **文本**产物没有 PDF 文件，不缓存。``put_cache``
        自带降级：缓存不可写返回原路径，不影响预览。

        Parameters
        ----------
        file_info : dict
            与 ``convert()`` 一致的文件信息。
        result : ConversionResult
            后端已产生的转换结果。

        Returns
        -------
        ConversionResult
            写入缓存后返回缓存内路径结果；非 PDF 产物或缓存失败时返回
            *result* 本身（转换流程照常）。
        """
        is_pdf_artifact = (
            result.content_type == "pdf"
            and isinstance(result.content, Path)
            and result.backend_used in ("libreoffice", "com")
        )
        if not is_pdf_artifact:
            return result
        try:
            from freeassetfilter.services.office_cache import put_cache
        except ImportError:
            return result
        try:
            cached_path = put_cache(file_info, result.content)
        except Exception:
            return result
        return ConversionResult(
            content_type=result.content_type,
            content=cached_path,
            backend_used=result.backend_used,
            truncated=result.truncated,
            message=result.message,
        )

    # ── 内部工具 ─────────────────────────────────────────────────────

    @staticmethod
    def _normalize_suffix(file_info: dict) -> str:
        """
        从 *file_info* 提取归一化后缀（小写、无前导点）。

        Parameters
        ----------
        file_info : dict
            含 ``"suffix"`` 键的文件信息；非字典或缺少键 / 值为 ``None``
            时返回 ``''``（保证 ``convert()`` 永不抛出）。

        Returns
        -------
        str
            归一化后缀（如 ``"docx"``），无法识别时返回 ``''``。
        """
        if not isinstance(file_info, dict):
            return ''
        suffix = file_info.get("suffix", '')
        if suffix is None:
            return ''
        return str(suffix).lower().lstrip(".")


__all__ = [
    "ConversionResult",
    "ERROR_MESSAGE",
    "LEGACY_SUFFIXES",
    "MODERN_SUFFIXES",
    "OfficeConverter",
    "SOFFICE_CONVERSION_TIMEOUT",
    "SUPPORTED_SUFFIXES",
    "clear_active_lo_popen",
    "get_active_lo_popen",
]
